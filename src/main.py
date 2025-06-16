import os
import re
import base64
from io import BytesIO
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PASTA_PPTS    = "caminho/para/a/pasta/com/arquivos_pptx"   # caminho para os arquivos .pptx com imagens
ARQUIVO_SAIDA = "caminho/para/salvar/relatorio.xlsx"       # caminho de saída do arquivo gerado
ALTURA_PADRAO = 234                                         # altura em pixels para redimensionamento da imagem (ajuste conforme visual do BI)
PREFIXO       = "PREFIXO_DO_ARQUIVO"                        # prefixo dos arquivos que devem ser processados (ex: "RELATORIO_SEMANAL")

# Palavras para sub-localização (áreas)
SUB_LOC_KEYWORDS = [
    'almoxarifado', 'vestiário', 'jardim', 'refeitório', 'corredor',
    'banheiro', 'elevador', 'garagem', 'auditório', 'área de serviço', 'recepção'
]

# Mapeamento de palavras-chave para categoria
CATEGORY_MAP = {
    'fachada': 'Fachada',
    'prumada': 'Prumada',
    'spk': 'SPCI',
    'elétrica': 'Instalações Elétricas',
    'ventilação': 'HVAC',
    'piso': 'Civil', 'divisória': 'Civil', 'drywall': 'Civil', 'pintura': 'Civil',
    'sdai': 'SDAI',
    'hidráulica': 'Instalações Hidráulicas', 'hidrossanitária': 'Instalações Hidráulicas'
}

# ————————— Funções de extração —————————————————

def _extrair_data(texto: str) -> str | None:
    m = re.search(r"\b(\d{1,2})[./-](\d{1,2})(?:[./-](\d{2,4}))?\b", texto)
    if not m:
        return None
    d, M, a = m.group(1), m.group(2), m.group(3)
    if a:
        a = f"20{a}" if len(a)==2 else a
        return f"{int(d):02d}/{int(M):02d}/{a}"
    return f"{int(d):02d}/{int(M):02d}"


def _extrair_pavimento(texto: str) -> str | None:
    txt = texto.lower()
    # 1) gedise: "1º gedise" → "1SS-G"
    m = re.search(r'\b(\d+)(?:º|o)?\s*gedise\b', txt)
    if m:
        return f"{m.group(1)}SS-G"

    # 2) subsolo
    m = re.search(r'\b(\d+)\s*sub\s*solo\b', txt)
    if m:
        return f"{m.group(1)}SS"

    # 3) número→pavimento/andar
    m = re.search(r'\b(\d+)[º°]?\s*(?:pavimento|andar)\b', txt)
    if m:
        return f"{m.group(1)}º Andar"

    # 4) pavimento/andar→número
    m = re.search(r'\b(?:pavimento|andar)\s*(\d+)\b', txt)
    if m:
        return f"{m.group(1)}º Andar"

    # 5) térreo isolado
    if re.search(r'\bt[eé]rreo\b', txt):
        return "Térreo"

    return None


def _extrair_servico(texto: str) -> str | None:
    # 1) junta linhas e espaços
    txt = " ".join(texto.split())

    # 2) regex robusta que inclui gedise como keyword
    pattern = re.compile(r'''
        ^(?:Dia\s*)?                              # "Dia" opcional
        \d{1,2}[./-]\d{1,2}(?:[./-]\d{2,4})?      # data
        \s*[-–—]?\s*
        (?:Gedise|Sub\s*solo|Pavimento|Andar)     # keyword (inclui Gedise)
        \s*(?:\d+[º°]?|t[eé]rreo)?                # número ou "Térreo"
        \s*[-–—]?\s*
        (.*)                                      # ← group1 = descrição do serviço
    ''', re.IGNORECASE | re.VERBOSE)

    m = pattern.match(txt)
    if m:
        return m.group(1).strip() or None

    # 3) fallback genérico: remove GEDISE, data e pavimento
    serv = re.sub(r'\bDia\s*', '', txt, flags=re.IGNORECASE)
    serv = re.sub(r'\b\d{1,2}[./-]\d{1,2}(?:[./-]\d{2,4})?\b', '', serv)
    serv = re.sub(r'\b(?:Gedise|Pavimento|Andar)\s*(?:\d+[º°]?|t[eé]rreo)\b', '',
                  serv, flags=re.IGNORECASE)
    serv = re.sub(r'\b\d+\s*sub\s*solo\b', '', serv, flags=re.IGNORECASE)
    # remove hífen inicial e espaços sobrando
    serv = re.sub(r'^[\s\-–—]+', '', serv).strip()

    return serv or None


def extrair_sub_localizacao(servico: str) -> str | None:
    txt = servico.lower()
    for kw in SUB_LOC_KEYWORDS:
        if kw in txt:
            return kw.title()
    return None


def extrair_categoria(servico: str) -> str | None:
    txt = servico.lower()
    for key, cat in CATEGORY_MAP.items():
        if key in txt:
            return cat
    return None


def encontrar_legenda_por_caixa(picture, shapes):
    base_y = picture.top + picture.height
    cand = []
    for shp in shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and shp.has_text_frame:
            if shp.top >= base_y and shp.left < picture.left + picture.width and shp.left + shp.width > picture.left:
                txt = shp.text_frame.text.strip()
                if txt:
                    cand.append((shp.top - base_y, txt))
    if not cand:
        return ""
    return sorted(cand, key=lambda x: x[0])[0][1]


def converter_para_base64(blob: bytes) -> str:
    with Image.open(BytesIO(blob)) as img:
        h = ALTURA_PADRAO / img.height
        w = int(img.width * h)
        r = img.resize((w, ALTURA_PADRAO), Image.LANCZOS)
        buf = BytesIO(); fmt = img.format or "PNG"
        r.save(buf, format=fmt)
        return f"data:image/{fmt.lower()};base64,{base64.b64encode(buf.getvalue()).decode()}"


def main():
    registros = []
    for nome in os.listdir(PASTA_PPTS):
        if not nome.startswith(PREFIXO) or not nome.lower().endswith(".pptx"):
            continue
        prs = Presentation(os.path.join(PASTA_PPTS, nome))
        total = len(prs.slides)
        for i, slide in enumerate(prs.slides, start=1):
            if i in (1, total):
                continue
            idx = 0
            for shp in slide.shapes:
                if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    idx += 1
                    legend = encontrar_legenda_por_caixa(shp, slide.shapes)
                    data   = _extrair_data(legend) or ''
                    pav    = _extrair_pavimento(legend) or ''
                    serv   = _extrair_servico(legend) or ''
                    registros.append({
                        'ArquivoPPTX': nome,
                        'ImagemIndex': idx,
                        'Legenda': legend,
                        'DataDetectada': data,
                        'Pavimento': pav,
                        'Servico': serv,
                        'SubLocalizacao': extrair_sub_localizacao(serv) or '',
                        'Categoria': extrair_categoria(serv) or '',
                        'Base64': converter_para_base64(shp.image.blob)
                    })
    df = pd.DataFrame(registros, columns=[
        'ArquivoPPTX', 'ImagemIndex', 'Legenda', 'DataDetectada',
        'Pavimento', 'Servico', 'SubLocalizacao', 'Categoria', 'Base64'
    ])
    df.to_excel(ARQUIVO_SAIDA, index=False)
    print(f"{len(registros)} imagens processadas. Saída: {ARQUIVO_SAIDA}")

if __name__ == '__main__':
    main()
